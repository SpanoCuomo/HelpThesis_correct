# -*- coding: utf-8 -*-
"""
Created on Wed Jul 17 13:56:33 2024

@author: UTENTE
"""

from pptx import Presentation
import os
import csv

def replace_text_in_pptx(file_path, replacements):
    # Carica la presentazione
    prs = Presentation(file_path)
    
    # Cicla su tutte le slide
    for slide in prs.slides:
        # Cicla su tutti gli elementi di testo in una slide
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements:
                        if old_text in run.text:
                            if '\\item' in new_text:
                                # Gestisce il caso in cui new_text contiene elementi puntati
                                bullet_points = new_text.split('\\item ')
                                paragraph.clear()  # Cancella il contenuto esistente del paragrafo
                                for point in bullet_points:
                                    if point.strip():
                                        p = paragraph.add_run()
                                        p.text = f"- {point.strip()}"
                            else:
                                run.text = run.text.replace(old_text, new_text)
    
    # Salva la presentazione modificata
    prs.save('modified_presentation.pptx')
    print("Text replacement complete and file saved as 'modified_presentation.pptx'")

def read_replacements(file_path):
    replacements = []
    with open(file_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file, quotechar='"', delimiter=',', quoting=csv.QUOTE_ALL, skipinitialspace=True)
        for idx, row in enumerate(reader):
            if len(row) == 2:
                old_text = row[0].strip().strip('"')
                new_text = row[1].strip().strip('"').replace('\\n', '\n')  # Sostituisci \n con il carattere di nuova riga
                replacements.append((old_text, new_text))
            else:
                print(f"Skipping invalid line {idx + 1}: {','.join(row)}")
    return replacements

# Leggi le sostituzioni dal file txt
replacements = read_replacements(r'C:/Users/UTENTE/Downloads/replacements_2.txt')

# Verifica se il file PowerPoint esiste nel percorso specificato
pptx_path = 'C:/Users/UTENTE/Desktop/Cavicchia_Copy.pptx'
if not os.path.exists(pptx_path):
    print(f"File not found: {pptx_path}")
else:
    # Usa questa funzione per sostituire il testo nel tuo file PowerPoint
    replace_text_in_pptx(pptx_path, replacements)
